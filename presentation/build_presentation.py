import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK = "0A0F1A"
NAVY = "003580"
NAVY_MID = "004099"
NAVY_DIM = "002060"
NAVY_TINT = "E8EFF8"
MIST = "F4F5F7"
BORDER = "D8DCE3"
MID = "6B7280"
WHITE = "FFFFFF"
ALERT = "B91C1C"

FONT_DISPLAY = "Plus Jakarta Sans"
FONT_BODY = "Inter"

EMU_W = 13.333
EMU_H = 7.5

prs = Presentation()
prs.slide_width = Inches(EMU_W)
prs.slide_height = Inches(EMU_H)
BLANK = prs.slide_layouts[6]


def c(hexstr):
    return RGBColor.from_string(hexstr)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = c(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = c(line)
        sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def write(tf, paras):
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(para.get("space_after", 0))
        p.space_before = Pt(para.get("space_before", 0))
        p.line_spacing = para.get("line_spacing", 1.0)
        for r in para["runs"]:
            run = p.add_run()
            run.text = r["text"]
            f = run.font
            f.name = r.get("font", FONT_BODY)
            f.size = Pt(r.get("size", 14))
            f.bold = r.get("bold", False)
            f.italic = r.get("italic", False)
            f.color.rgb = c(r.get("color", INK))
    return tf


def header(slide, kicker, title, page_no=None):
    write(textbox(slide, 0.7, 0.5, 11.9, 0.35), [{"runs": [{"text": kicker, "size": 12, "bold": True, "color": NAVY}]}])
    write(textbox(slide, 0.7, 0.82, 11.9, 0.75), [{"runs": [{"text": title, "size": 28, "bold": True, "color": INK, "font": FONT_DISPLAY}]}])
    rect(slide, 0.7, 1.72, 1.1, 0.055, fill=NAVY)
    footer(slide, page_no)


def footer(slide, page_no):
    rect(slide, 0.7, 7.08, 11.93, 0.015, fill=BORDER)
    write(textbox(slide, 0.7, 7.14, 9.0, 0.3), [{"runs": [{"text": "Abdul Latif · BINUS University — Laporan Magang Diskominfotik DKI Jakarta", "size": 9, "color": MID}]}])
    write(textbox(slide, 11.8, 7.14, 0.83, 0.3), [{"align": PP_ALIGN.RIGHT, "runs": [{"text": page_no, "size": 9, "color": MID}]}])


def number_circle(slide, cx, cy, d, num, fill=NAVY, text_color=WHITE, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d))
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = c(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = c(line)
        sp.line.width = Pt(1.5)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.name = FONT_DISPLAY
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = c(text_color)


def card(slide, x, y, w, h, accent, title, items, sub=None, item_size=12, title_size=17, bullets_from=3.05):
    rect(slide, x, y, w, h, fill=MIST, line=BORDER, line_w=1.0, radius=0.045, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, x, y, w, 0.09, fill=accent)
    write(textbox(slide, x + 0.28, y + 0.32, w - 0.56, 0.5), [{"runs": [{"text": title, "size": title_size, "bold": True, "color": INK, "font": FONT_DISPLAY}]}])
    paras = []
    yy = y + 0.95
    if sub:
        write(textbox(slide, x + 0.28, yy, w - 0.56, 0.42), [{"runs": [{"text": sub, "size": 10.5, "italic": True, "color": MID}]}])
        yy += 0.48
    for it in items:
        paras.append({"space_after": 8, "line_spacing": 1.2, "runs": [
            {"text": "\u25AA  ", "size": item_size, "bold": True, "color": accent},
            {"text": it, "size": item_size, "color": INK},
        ]})
    write(textbox(slide, x + 0.28, yy, w - 0.56, y + h - yy - 0.1), paras)
    return yy


# ------------------------------------------------------------------ SLIDE 1
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
rect(s, 0, 6.6, EMU_W, 0.9, fill=NAVY_DIM)

rect(s, 0.7, 0.75, 1.9, 0.5, fill=NAVY_TINT, radius=0.18, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
write(textbox(s, 0.7, 0.75, 1.9, 0.5, anchor=MSO_ANCHOR.MIDDLE), [{
    "align": PP_ALIGN.CENTER,
    "runs": [{"text": "CSIRT DKI JAKARTA", "size": 11, "bold": True, "color": NAVY}],
}])

rect(s, 11.43, 0.75, 1.2, 0.5, fill=WHITE, radius=0.18, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
write(textbox(s, 11.43, 0.75, 1.2, 0.5, anchor=MSO_ANCHOR.MIDDLE), [{
    "align": PP_ALIGN.CENTER,
    "runs": [{"text": "LIVE STAGING", "size": 10, "bold": True, "color": NAVY}],
}])

write(textbox(s, 0.7, 1.65, 11.9, 0.4), [{"runs": [{"text": "LAPORAN MAGANG · DISKOMINFOTIK DKI JAKARTA", "size": 13, "bold": True, "color": NAVY_TINT}]}])
write(textbox(s, 0.7, 2.05, 11.9, 2.1), [
    {"space_after": 10, "runs": [{"text": "Modernisasi Portal CSIRT DKI Jakarta", "size": 44, "bold": True, "color": WHITE, "font": FONT_DISPLAY}]},
    {"space_after": 6, "line_spacing": 1.25, "runs": [{"text": "Redesain UI/UX agar Mudah Digunakan Masyarakat Umum", "size": 17, "bold": True, "color": NAVY_TINT}]},
    {"runs": [{"text": "Aksesibilitas WCAG 2.1 · Dibangun Ulang dengan Laravel 12", "size": 15, "color": NAVY_TINT}]},
])
rect(s, 0.72, 4.25, 2.0, 0.07, fill=WHITE)

write(textbox(s, 0.7, 4.65, 11.9, 1.6), [
    {"space_after": 6, "runs": [{"text": "Abdul Latif", "size": 20, "bold": True, "color": WHITE}]},
    {"space_after": 4, "runs": [{"text": "Mahasiswa Computer Science — BINUS University", "size": 13, "color": NAVY_TINT}]},
    {"runs": [{"text": "Periode Magang: Februari – Agustus 2026   ·   Pembimbing Lapangan: Pak Andy & Bu Rina", "size": 13, "color": NAVY_TINT}]},
])

write(textbox(s, 0.7, 6.7, 11.9, 0.7, anchor=MSO_ANCHOR.MIDDLE), [{"runs": [{"text": "csirt.jakarta.go.id — Computer Security Incident Response Team Pemerintah Provinsi DKI Jakarta", "size": 11, "color": NAVY_TINT}]}])

# ------------------------------------------------------------------ SLIDE 2
s = prs.slides.add_slide(BLANK)
header(s, "PROFIL INSTANSI", "Diskominfotik DKI Jakarta & JakartaProv-CSIRT", "2 / 12")

cw = (11.93 - 0.35) / 2
lx = 0.7
rx = lx + cw + 0.35
cy = 2.1
ch = 3.85

for cx, logo, title, items in [
    (lx, r"K:\GitHub\new-csirt\presentation\logo_diskominfo.png", "Diskominfotik DKI Jakarta", [
        "Perangkat Daerah urusan Komunikasi & Informatika, Statistik, dan Persandian di wilayah DKI Jakarta",
        "Dipimpin Kepala Dinas — bertanggung jawab kepada Gubernur melalui Sekretaris Daerah",
        "Menaungi Bidang Siber, Sandi & Aplikasi — rumah bagi JakartaProv-CSIRT",
    ]),
    (rx, r"K:\GitHub\new-csirt\presentation\csirt-main-logo.png", "JakartaProv-CSIRT", [
        "Tim Tanggap Insiden Siber Pemprov DKI Jakarta — ditetapkan SK Sekda (Pj.) No. 41 Tahun 2020",
        "Mitigasi, koordinasi, penanggulangan & pemulihan insiden siber sektor Pemda DKI Jakarta",
        "Layanan: lapor insiden, peringatan keamanan, berita & panduan keamanan siber",
        "Hotline 0813-8887-0152 · csirt@jakarta.go.id · Balaikota Blok H Lt. 13",
    ]),
]:
    rect(s, cx, cy, cw, ch, fill=MIST, line=BORDER, line_w=1.0, radius=0.045, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, cy, cw, 0.09, fill=NAVY)
    s.shapes.add_picture(logo, Inches(cx + 0.28), Inches(cy + 0.34), height=Inches(0.85))
    write(textbox(s, cx + 1.25, cy + 0.4, cw - 1.53, 0.7, anchor=MSO_ANCHOR.MIDDLE), [{
        "runs": [{"text": title, "size": 16, "bold": True, "color": INK, "font": FONT_DISPLAY}],
    }])
    paras = []
    for it in items:
        paras.append({"space_after": 8, "line_spacing": 1.2, "runs": [
            {"text": "\u25AA  ", "size": 12, "bold": True, "color": NAVY},
            {"text": it, "size": 12, "color": INK},
        ]})
    write(textbox(s, cx + 0.28, cy + 1.45, cw - 0.56, ch - 1.55), paras)

rect(s, 0.7, 6.25, 11.93, 0.7, fill=MIST, line=BORDER, line_w=1.0, radius=0.12, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.7, 6.25, 0.1, 0.7, fill=NAVY)
write(textbox(s, 1.05, 6.25, 11.3, 0.7, anchor=MSO_ANCHOR.MIDDLE), [{
    "line_spacing": 1.2,
    "runs": [
        {"text": "Lokasi magang: ", "size": 13, "bold": True, "color": NAVY},
        {"text": "tim pengembang portal resmi csirt.jakarta.go.id — Februari s.d. Agustus 2026.", "size": 13, "color": INK},
    ],
}])

# ------------------------------------------------------------------ SLIDE 3
s = prs.slides.add_slide(BLANK)
header(s, "AWAL MULA PROYEK", "Usulan Februari: UI Audit & Metode", "3 / 12")

cw = (11.93 - 0.35) / 2
lx = 0.7
rx = lx + cw + 0.35

card(s, lx, 2.1, cw, 3.9, ALERT, "Temuan UI Audit (Feb 2026)",
     sub="Berdasarkan demonstrasi landing page lama",
     items=[
        "Tampilan lama terlihat rapi, namun banyak elemen yang bisa diklik tidak menuju ke mana pun",
        "Footer tertutup elemen lain, terlalu besar, dan penuh informasi",
        "Situs terlalu fokus pada kebutuhan ahli IT — membingungkan warga biasa",
     ], item_size=12.5, title_size=15, bullets_from=3.1)

card(s, rx, 2.1, cw, 3.9, NAVY, "Metode — Idea Analysis",
     sub="Empat langkah dari proposal Februari",
     items=[
        "UI Audit — memetakan ulang elemen pada halaman yang tersedia",
        "Redesign — desain ulang antarmuka dengan Figma",
        "Front-end Cleanup — implementasi ke codebase website",
        "Survey & Reiterate — kritik & saran pengguna umum, lalu sesuaikan",
     ], item_size=12.5, title_size=15, bullets_from=3.1)

rect(s, 0.7, 6.25, 11.93, 0.7, fill=MIST, line=BORDER, line_w=1.0, radius=0.12, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.7, 6.25, 0.1, 0.7, fill=NAVY)
write(textbox(s, 1.05, 6.25, 11.3, 0.7, anchor=MSO_ANCHOR.MIDDLE), [{
    "line_spacing": 1.2,
    "runs": [
        {"text": "Dari proposal Februari → proyek modernisasi penuh. ", "size": 13, "bold": True, "color": NAVY},
        {"text": "Langkah 'Survey & Reiterate' berlanjut ke usability testing — fokus skripsi.", "size": 13, "color": INK},
    ],
}])

# ------------------------------------------------------------------ SLIDE 4
s = prs.slides.add_slide(BLANK)
header(s, "DATA PENGGUNA", "Survei Pengguna: Bukti Kebutuhan Redesain", "4 / 12")

stats = [
    ("50%", "terganggu warna & desain"),
    ("44%", "keluhan ukuran teks/font"),
    ("31%", "tata letak / navigasi / loading"),
    ("13%", "error saat submit laporan"),
]
sw = (11.93 - 0.9) / 4
sx = 0.7
for val, label in stats:
    rect(s, sx, 2.05, sw, 0.95, fill=NAVY, radius=0.09, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    write(textbox(s, sx + 0.2, 2.05, sw - 0.4, 0.95, anchor=MSO_ANCHOR.MIDDLE), [
        {"align": PP_ALIGN.CENTER, "runs": [{"text": val, "size": 24, "bold": True, "color": WHITE, "font": FONT_DISPLAY}]},
        {"align": PP_ALIGN.CENTER, "runs": [{"text": label, "size": 10.5, "color": NAVY_TINT}]},
    ])
    sx += sw + 0.3

cw = (11.93 - 0.35) / 2
lx = 0.7
rx = lx + cw + 0.35

card(s, lx, 3.35, cw, 3.1, ALERT, "Keluhan Terbanyak (N = 16)",
     sub="Survei pengguna umum, 4–5 Februari 2026",
     items=[
        "Warna & desain — diminta biru navy, bukan oranye terang",
        "Ukuran teks/font yang mengganggu",
        "Penempatan fitur, navigasi & waktu loading",
        "Gagal mengirim laporan insiden (error database)",
     ])

card(s, rx, 3.35, cw, 3.1, NAVY, "Jawaban Melalui Redesain",
     sub="Setiap keluhan dipetakan ke solusi",
     items=[
        "Design system navy #003580 yang konsisten",
        "Widget aksesibilitas: teks 80–150%, kontras tinggi",
        "Struktur layout dipetakan ulang (UI audit)",
        "Wizard 3 langkah + penanganan error + rate limit",
     ])

# ------------------------------------------------------------------ SLIDE 5
s = prs.slides.add_slide(BLANK)
header(s, "DESKRIPSI PROYEK", "Permasalahan & Tujuan", "5 / 12")

cards = [
    {
        "title": "Latar Belakang",
        "accent": NAVY,
        "body": [
            {"kind": "text", "text": "Portal CSIRT DKI Jakarta adalah kanal publik untuk melaporkan insiden siber dan literasi keamanan digital bagi warga serta pegawai pemerintah."},
            {"kind": "text", "text": "Versi lama (Yii) terlalu fokus pada kebutuhan ahli IT — membingungkan warga biasa yang hanya ingin melapor."},
            {"kind": "text", "text": "Survei Februari: 50% terganggu warna & desain, 13% gagal saat submit laporan."},
        ],
    },
    {
        "title": "Permasalahan",
        "accent": ALERT,
        "body": [
            {"kind": "bullet", "text": "Portal expert-centric, penuh istilah teknis"},
            {"kind": "bullet", "text": "Form laporan 1 halaman panjang, kotak kosong & skor CVSS"},
            {"kind": "bullet", "text": "Tidak ramah gangguan penglihatan / buta warna"},
            {"kind": "bullet", "text": "Konten hanya bisa diubah lewat kode program"},
        ],
    },
    {
        "title": "Tujuan",
        "accent": NAVY_MID,
        "body": [
            {"kind": "bullet", "text": "Melapor insiden mudah bagi masyarakat umum"},
            {"kind": "bullet", "text": "Wizard 3 langkah yang menuntun pengguna"},
            {"kind": "bullet", "text": "Aksesibilitas standar layanan publik (WCAG 2.1)"},
            {"kind": "bullet", "text": "Staf mengelola konten tanpa menyentuh kode"},
        ],
    },
]

cw = (11.93 - 0.7) / 3
x = 0.7
for card_data in cards:
    rect(s, x, 2.1, cw, 4.35, fill=MIST, line=BORDER, line_w=1.0, radius=0.045, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, 2.1, cw, 0.09, fill=card_data["accent"])
    write(textbox(s, x + 0.28, 2.42, cw - 0.56, 0.5), [{"runs": [{"text": card_data["title"], "size": 17, "bold": True, "color": INK, "font": FONT_DISPLAY}]}])
    tf = textbox(s, x + 0.28, 3.05, cw - 0.56, 3.3)
    paras = []
    for item in card_data["body"]:
        if item["kind"] == "text":
            paras.append({"space_after": 10, "line_spacing": 1.25, "runs": [{"text": item["text"], "size": 12, "color": INK}]})
        else:
            paras.append({"space_after": 8, "line_spacing": 1.2, "runs": [
                {"text": "\u25AA  ", "size": 12, "bold": True, "color": card_data["accent"]},
                {"text": item["text"], "size": 12, "color": INK},
            ]})
    write(tf, paras)
    x += cw + 0.35

# ------------------------------------------------------------------ SLIDE 6
s = prs.slides.add_slide(BLANK)
header(s, "TAHAPAN PENGERJAAN", "Progress Pengerjaan", "6 / 12")

steps = [
    ("1", "Maret", "Struktur & Database", "Struktur dasar Laravel + tabel laporan insiden & berita", NAVY),
    ("2", "Mar – Apr", "Form Laporan", "Wizard 3 langkah + halaman Terima Kasih", NAVY),
    ("3", "Apr – Mei", "Admin & Konten", "Login & dashboard aman untuk staf internal", NAVY),
    ("4", "Mei – sekarang", "Desain & Aksesibilitas", "Galeri infografis + widget aksesibilitas", NAVY),
    ("5", "Berikutnya", "Usability Testing", "Bandingkan lama vs baru, ukur error & waktu", NAVY_MID),
]

col_w = 11.93 / 5
cx0 = 0.7
cy = 2.6
rect(s, cx0 + 0.6, cy - 0.045, col_w * 4 - 1.2, 0.09, fill=BORDER)
for i, (num, period, title, desc, accent) in enumerate(steps):
    cx = cx0 + col_w * i + col_w / 2
    future = (i == 4)
    if future:
        number_circle(s, cx, cy, 0.92, num, fill=NAVY_TINT, text_color=NAVY_MID, line=NAVY_MID)
    else:
        number_circle(s, cx, cy, 0.92, num, fill=accent)
    write(textbox(s, cx - col_w / 2 + 0.18, cy + 0.62, col_w - 0.36, 0.4), [{
        "align": PP_ALIGN.CENTER,
        "runs": [{"text": period.upper(), "size": 10.5, "bold": True, "color": MID}],
    }])
    write(textbox(s, cx - col_w / 2 + 0.18, cy + 1.0, col_w - 0.36, 0.5), [{
        "align": PP_ALIGN.CENTER,
        "runs": [{"text": title, "size": 14.5, "bold": True, "color": INK, "font": FONT_DISPLAY}],
    }])
    write(textbox(s, cx - col_w / 2 + 0.18, cy + 1.52, col_w - 0.36, 1.4), [{
        "align": PP_ALIGN.CENTER,
        "line_spacing": 1.2,
        "runs": [{"text": desc, "size": 11, "color": MID}],
    }])

rect(s, 0.7, 5.75, 11.93, 0.85, fill=MIST, line=BORDER, line_w=1.0, radius=0.12, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.7, 5.75, 0.1, 0.85, fill=NAVY)
write(textbox(s, 1.05, 5.75, 11.3, 0.85, anchor=MSO_ANCHOR.MIDDLE), [{
    "line_spacing": 1.2,
    "runs": [
        {"text": "Pendekatan iteratif: ", "size": 13, "bold": True, "color": NAVY},
        {"text": "setiap fitur melalui siklus rancang → bangun → uji → perbaiki, dengan umpan balik dari pembimbing lapangan.", "size": 13, "color": INK},
    ],
}])

# ------------------------------------------------------------------ SLIDE 7
s = prs.slides.add_slide(BLANK)
header(s, "OUTPUT", "Sebelum vs Sesudah: Form Laporan Insiden", "7 / 12")

cw = (11.93 - 0.35) / 2
lx = 0.7
rx = lx + cw + 0.35

card(s, lx, 2.1, cw, 4.45, ALERT, "SEBELUM — Portal Lama (Yii)", [
    "Satu halaman sangat panjang, penuh kotak teks kosong",
    "Meminta data teknis rumit — skor CVSS",
    "Istilah yang tidak dipahami warga awam",
    "Warga bingung — banyak yang tidak jadi melapor",
], sub=None, item_size=12.5, title_size=15, bullets_from=3.1)

card(s, rx, 2.1, cw, 4.45, NAVY, "SESUDAH — Portal Baru (Laravel)", [
    "Wizard 3 langkah yang menuntun pengguna:",
    "Langkah 1 — Siapa yang melapor?",
    "Langkah 2 — Website apa yang diserang?",
    "Langkah 3 — Apa yang terjadi?",
    "Istilah teknis dibuat opsional — melapor cepat & mudah",
], sub=None, item_size=12.5, title_size=15, bullets_from=3.1)

arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(lx + cw - 0.05), Inches(3.4), Inches(0.45), Inches(0.55))
arrow.shadow.inherit = False
arrow.fill.solid()
arrow.fill.fore_color.rgb = c(NAVY)
arrow.line.fill.background()

# ------------------------------------------------------------------ SLIDE 8
s = prs.slides.add_slide(BLANK)
header(s, "OUTPUT", "Sebelum vs Sesudah: Aksesibilitas & Konten", "8 / 12")

stats = [
    ("9", "halaman publik didesain ulang"),
    ("6", "modul panel admin"),
    ("3", "mode kontras"),
    ("60", "permintaan/menit rate limit"),
]
sw = (11.93 - 0.9) / 4
sx = 0.7
for val, label in stats:
    rect(s, sx, 2.05, sw, 0.95, fill=NAVY, radius=0.09, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    write(textbox(s, sx + 0.2, 2.05, sw - 0.4, 0.95, anchor=MSO_ANCHOR.MIDDLE), [
        {"align": PP_ALIGN.CENTER, "runs": [{"text": val, "size": 24, "bold": True, "color": WHITE, "font": FONT_DISPLAY}]},
        {"align": PP_ALIGN.CENTER, "runs": [{"text": label, "size": 10.5, "color": NAVY_TINT}]},
    ])
    sx += sw + 0.3

cw = (11.93 - 0.35) / 2
lx = 0.7
rx = lx + cw + 0.35

card(s, lx, 3.35, cw, 3.1, NAVY, "Widget Aksesibilitas",
     sub="Sebelum: tidak ramah gangguan penglihatan / buta warna",
     items=[
        "Ukuran teks 80% – 150%",
        "Spasi baris & spasi huruf",
        "Kontras tinggi & dark mode",
        "Preferensi tersimpan di browser",
     ])

card(s, rx, 3.35, cw, 3.1, NAVY_MID, "Admin — Manajemen Konten",
     sub="Sebelum: konten hanya bisa diubah lewat kode",
     items=[
        "Login & dashboard aman untuk staf internal",
        "Tambah / ubah / hapus berita, peringatan, infografis",
        "Galeri infografis rapi & cepat dimuat",
        "Pagination & tab untuk 6 jenis konten",
     ])

# ------------------------------------------------------------------ SLIDE 9
s = prs.slides.add_slide(BLANK)
header(s, "PERJALANAN MAGANG", "Perjalanan Magang Bulan 1–5", "9 / 12")

months = [
    ("BULAN 1", "Audit & Proposal", "Usulan redesign UI; memetakan masalah pengguna sejak Februari", NAVY),
    ("BULAN 2", "Front-End & MVC", "HTML/CSS/Bootstrap; arsitektur MVC Laravel; prinsip WCAG 2.1", NAVY),
    ("BULAN 3", "CMS Admin (CRUD)", "CRUD 6 modul konten; dari perancang UI → full-stack developer", NAVY),
    ("BULAN 4", "Aksesibilitas", "Widget teks 80–150% & kontras; optimasi query; galeri lightbox", NAVY),
    ("BULAN 5", "Deploy & Audit", "Deployment staging Railway.app; audit arsitektur & dokumentasi", NAVY_MID),
]
col_w = 11.93 / 5
for i, (kicker, title, desc, accent) in enumerate(months):
    cx = 0.7 + col_w * i
    w = col_w - 0.24
    rect(s, cx, 2.1, w, 3.3, fill=MIST, line=BORDER, line_w=1.0, radius=0.06, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, 2.1, w, 0.09, fill=accent)
    write(textbox(s, cx + 0.15, 2.42, w - 0.3, 0.3), [{"runs": [{"text": kicker, "size": 10.5, "bold": True, "color": MID}]}])
    write(textbox(s, cx + 0.15, 2.72, w - 0.3, 0.6), [{"runs": [{"text": title, "size": 14, "bold": True, "color": INK, "font": FONT_DISPLAY}]}])
    write(textbox(s, cx + 0.15, 3.42, w - 0.3, 1.8), [{"line_spacing": 1.2, "runs": [{"text": desc, "size": 10.5, "color": MID}]}])

rect(s, 0.7, 5.65, 11.93, 1.25, fill=MIST, line=BORDER, line_w=1.0, radius=0.12, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.7, 5.65, 0.1, 1.25, fill=NAVY)
write(textbox(s, 1.05, 5.78, 11.3, 0.3), [{"runs": [{"text": "3 atribut lulusan BINUS yang terasah selama magang:", "size": 12, "bold": True, "color": NAVY}]}])
attrs = [
    ("Critical & Creative Thinking", "Jonassen (2000)"),
    ("Digital & Technology Fluency", "TAM — Davis (1989)"),
    ("Applied Management Skills", "Self-Efficacy — Bandura (1997)"),
]
ax = 1.05
for name, cite in attrs:
    write(textbox(s, ax, 6.14, 3.66, 0.66), [
        {"space_after": 2, "runs": [{"text": name, "size": 11.5, "bold": True, "color": INK}]},
        {"runs": [{"text": cite, "size": 10, "italic": True, "color": MID}]},
    ])
    ax += 3.82

# ------------------------------------------------------------------ SLIDE 10
s = prs.slides.add_slide(BLANK)
header(s, "REFLEKSI", "Pembelajaran (Lesson Learned)", "10 / 12")

cols = [
    {
        "title": "Teknis",
        "accent": NAVY,
        "items": [
            "Wizard 3 langkah mengalahkan form flat — pengguna menyerah pada halaman panjang",
            "Aksesibilitas (WCAG 2.1) bukan pelengkap — layanan publik wajib inklusif",
            "Struktur database yang benar sejak awal mempercepat fitur berikutnya",
            "Dokumentasi paralel menjaga konteks bagi pengembang berikutnya",
        ],
    },
    {
        "title": "Non-Teknis",
        "accent": NAVY_MID,
        "items": [
            "Empati terhadap pengguna awam — warga bukan ahli IT",
            "Koordinasi dengan pembimbing (Pak Andy & Bu Rina) di lingkungan pemerintah",
            "Manajemen waktu antar deliverable — testimoni (27 Juli) & pre-thesis (22 Agustus)",
            "Pengujian pada perangkat & koneksi nyata mengungkap masalah tersembunyi",
        ],
    },
]
cw = (11.93 - 0.35) / 2
x = 0.7
for col in cols:
    card(s, x, 2.1, cw, 4.35, col["accent"], col["title"], col["items"], sub=None)
    x += cw + 0.35

# ------------------------------------------------------------------ SLIDE 11
s = prs.slides.add_slide(BLANK)
header(s, "RENCANA", "Langkah Selanjutnya — Fokus Skripsi", "11 / 12")

rect(s, 0.7, 2.1, 11.93, 3.15, fill=MIST, line=BORDER, line_w=1.0, radius=0.045, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.7, 2.1, 11.93, 0.09, fill=NAVY)
write(textbox(s, 1.0, 2.42, 11.3, 0.5), [{"runs": [{"text": "Agar proyek siap diajukan sebagai topik skripsi di universitas:", "size": 15, "bold": True, "color": INK, "font": FONT_DISPLAY}]}])
tf = textbox(s, 1.0, 3.1, 11.3, 2.0)
write(tf, [
    {"space_after": 10, "line_spacing": 1.25, "runs": [
        {"text": "\u25AA  ", "size": 13, "bold": True, "color": NAVY},
        {"text": "Usability testing langsung dengan beberapa pengguna awam.", "size": 13, "color": INK},
    ]},
    {"space_after": 10, "line_spacing": 1.25, "runs": [
        {"text": "\u25AA  ", "size": 13, "bold": True, "color": NAVY},
        {"text": "Membandingkan website lama vs website baru.", "size": 13, "color": INK},
    ]},
    {"space_after": 10, "line_spacing": 1.25, "runs": [
        {"text": "\u25AA  ", "size": 13, "bold": True, "color": NAVY},
        {"text": "Membuktikan secara data: pengurangan error & waktu pengisian laporan.", "size": 13, "color": INK},
    ]},
    {"space_after": 0, "line_spacing": 1.25, "runs": [
        {"text": "\u25AA  ", "size": 13, "bold": True, "color": NAVY},
        {"text": "Menjadi fondasi proposal pre-thesis (22 Agustus 2026).", "size": 13, "color": INK},
    ]},
])

rect(s, 0.7, 5.6, 11.93, 0.9, fill=NAVY, radius=0.12, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
write(textbox(s, 0.7, 5.6, 11.93, 0.9, anchor=MSO_ANCHOR.MIDDLE), [{
    "align": PP_ALIGN.CENTER,
    "runs": [{"text": "Buktikan dengan data, bukan opini.", "size": 18, "bold": True, "color": WHITE, "font": FONT_DISPLAY}],
}])

# ------------------------------------------------------------------ SLIDE 12
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
rect(s, 0, 6.6, EMU_W, 0.9, fill=NAVY_DIM)

write(textbox(s, 0.7, 1.9, 11.9, 1.3), [{
    "align": PP_ALIGN.CENTER,
    "runs": [{"text": "Terima Kasih", "size": 48, "bold": True, "color": WHITE, "font": FONT_DISPLAY}],
}])
rect(s, 5.87, 3.2, 1.6, 0.07, fill=WHITE)
write(textbox(s, 0.7, 3.5, 11.9, 0.6), [{
    "align": PP_ALIGN.CENTER,
    "runs": [{"text": "Portal CSIRT DKI Jakarta yang lebih mudah digunakan masyarakat umum.", "size": 16, "color": NAVY_TINT}],
}])
rect(s, 4.22, 4.15, 4.9, 0.55, fill=NAVY_TINT, radius=0.18, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
write(textbox(s, 4.22, 4.15, 4.9, 0.55, anchor=MSO_ANCHOR.MIDDLE), [{
    "align": PP_ALIGN.CENTER,
    "runs": [{"text": "Demo: screenshot portal — menyusul saat live di Render", "size": 11, "bold": True, "color": NAVY}],
}])
write(textbox(s, 0.7, 5.05, 11.9, 0.6), [{
    "align": PP_ALIGN.CENTER,
    "runs": [{"text": "SESI TANYA JAWAB", "size": 14, "bold": True, "color": NAVY_TINT}],
}])
write(textbox(s, 0.7, 6.7, 11.9, 0.7, anchor=MSO_ANCHOR.MIDDLE), [{
    "align": PP_ALIGN.CENTER,
    "runs": [{"text": "Abdul Latif · BINUS University — csirt.jakarta.go.id", "size": 11, "color": NAVY_TINT}],
}])

OUT = r"K:\GitHub\new-csirt\presentation\Presentasi_Magang_CSIRT_DKI_Jakarta.pptx"
prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
